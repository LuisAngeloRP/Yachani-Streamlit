# pages/4_📤_upload.py
import streamlit as st
import os
import tempfile
from utils.document_manager import DocumentManager
from langchain_community.document_loaders import (
    PyPDFLoader, 
    UnstructuredWordDocumentLoader,
    UnstructuredEPubLoader,
    UnstructuredHTMLLoader,
    UnstructuredPowerPointLoader
)
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_openai.embeddings import OpenAIEmbeddings
from langchain_openai import ChatOpenAI
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import base64
from pathlib import Path
import re
import shutil

st.set_page_config(
    page_title="Subir Documento",
    page_icon="📤",
    layout="wide"
)

# Configuración de formatos soportados
SUPPORTED_FORMATS = {
    "pdf": ("PDF", ".pdf"),
    "docx": ("Word", ".docx"),
    "doc": ("Word", ".doc"),
    "epub": ("EPub", ".epub"),
    "txt": ("Text", ".txt"),
    "html": ("HTML", ".html"),
    "pptx": ("PowerPoint", ".pptx"),
    "ppt": ("PowerPoint", ".ppt")
}

def ensure_dir(path):
    """Asegura que un directorio exista."""
    os.makedirs(path, exist_ok=True)
    return path

def clean_filename(filename):
    """Limpia el nombre del archivo para que sea seguro."""
    return "".join(c if c.isalnum() or c in "._- " else "_" for c in filename)

def clean_text_with_ai(text: str, llm) -> str:
    """Usa IA para limpiar y estructurar mejor el texto."""
    try:
        prompt = f"""Por favor, limpia y estructura el siguiente texto manteniendo toda la información importante:
        1. Elimina caracteres extraños y formato innecesario
        2. Corrige errores obvios de formato
        3. Mantén la estructura de párrafos y secciones
        4. No agregues ni modifiques el contenido
        5. Asegura que el texto sea coherente y legible
        
        Texto: {text[:1500]}  # Limitamos para no usar muchos tokens
        """
        
        response = llm.invoke(prompt)
        return response.content
    except Exception as e:
        st.warning(f"No se pudo aplicar limpieza IA: {str(e)}")
        return text

def create_preview_image(file_path: str, output_path: str, file_type: str):
    """Crea una imagen de vista previa del documento."""
    try:
        if file_type == "pdf":
            doc = fitz.open(file_path)
            page = doc[0]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            pix.save(output_path)
            doc.close()
        elif file_type in ["ppt", "pptx"]:
            prs = Presentation(file_path)
            if len(prs.slides) > 0:
                slide = prs.slides[0]
                # Guardar la primera diapositiva como imagen
                # Nota: Esto requeriría una implementación adicional
                return False
        return True
    except Exception as e:
        st.warning(f"No se pudo crear vista previa: {str(e)}")
        return False

def get_document_loader(file_path: str, file_type: str):
    """Retorna el loader apropiado según el tipo de archivo."""
    loaders = {
        "pdf": PyPDFLoader,
        "docx": UnstructuredWordDocumentLoader,
        "doc": UnstructuredWordDocumentLoader,
        "epub": UnstructuredEPubLoader,
        "html": UnstructuredHTMLLoader,
        "txt": UnstructuredHTMLLoader,
        "pptx": UnstructuredPowerPointLoader,
        "ppt": UnstructuredPowerPointLoader
    }
    
    loader_class = loaders.get(file_type)
    if not loader_class:
        raise ValueError(f"Formato no soportado: {file_type}")
    
    return loader_class(file_path)

def create_download_link(file_path: str, link_text: str):
    """Crea un link de descarga para un archivo."""
    try:
        with open(file_path, "rb") as f:
            bytes_data = f.read()
        b64 = base64.b64encode(bytes_data).decode()
        filename = os.path.basename(file_path)
        mime_type = "application/octet-stream"
        href = f'<a href="data:{mime_type};base64,{b64}" download="{filename}" class="download-link">{link_text}</a>'
        return href
    except Exception as e:
        st.error(f"Error al crear link de descarga: {str(e)}")
        return None

def process_document(file, metadata, temp_dir):
    """Procesa el documento y crea el vectorstore."""
    try:
        # Determinar tipo de archivo
        file_extension = Path(file.name).suffix.lower()[1:]
        if file_extension not in SUPPORTED_FORMATS:
            return {"success": False, "error": "Formato de archivo no soportado"}

        # Preparar directorios
        safe_title = clean_filename(metadata["title"])
        doc_dir = ensure_dir(os.path.join("data", "processed_docs", safe_title))
        
        # Guardar archivo temporal
        temp_path = os.path.join(temp_dir, clean_filename(file.name))
        with open(temp_path, "wb") as f:
            f.write(file.getvalue())
        
        # Guardar copia del original
        original_path = os.path.join(doc_dir, f"original_{safe_title}{Path(file.name).suffix}")
        shutil.copy2(temp_path, original_path)
        
        # Crear vista previa
        preview_path = os.path.join(doc_dir, f"{safe_title}_preview.png")
        preview_created = create_preview_image(temp_path, preview_path, file_extension)
        
        # Procesar documento
        loader = get_document_loader(temp_path, file_extension)
        documents = loader.load()
        
        # Limpiar texto con IA (muestra)
        llm = ChatOpenAI(temperature=0, max_tokens=500)
        if documents:
            st.write("🔍 Analizando y limpiando el texto...")
            progress_bar = st.progress(0)
            
            # Procesar solo una muestra del texto
            sample_text = documents[0].page_content[:1500]
            cleaned_text = clean_text_with_ai(sample_text, llm)
            
            progress_bar.progress(50)
            st.info("✨ Muestra de texto limpiado (primer fragmento):")
            with st.expander("Ver muestra"):
                st.write(cleaned_text)
            
            progress_bar.progress(100)
        
        # Dividir en chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=150,
            separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""],
            length_function=len
        )
        chunks = text_splitter.split_documents(documents)
        
        # Crear vectorstore
        embeddings = OpenAIEmbeddings()
        vectorstore = Chroma.from_documents(
            documents=chunks,
            embedding=embeddings,
            persist_directory=doc_dir
        )
        
        return {
            "success": True,
            "num_pages": len(documents),
            "num_chunks": len(chunks),
            "vectorstore_path": doc_dir,
            "original_path": original_path,
            "preview_path": preview_path if preview_created else None,
            "file_type": file_extension,
            "file_size": os.path.getsize(original_path)
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": str(e)
        }

def main():
    
    st.title("📤 Subir Nuevo Documento")
    
    # Mostrar formatos soportados
    st.write("📌 Formatos soportados:")
    format_cols = st.columns(len(SUPPORTED_FORMATS))
    for col, (format_name, format_info) in zip(format_cols, SUPPORTED_FORMATS.items()):
        with col:
            st.markdown(f"**{format_info[0]}** ({format_info[1]})")
    
    doc_manager = DocumentManager()
    
    # Progress tracking
    if 'upload_step' not in st.session_state:
        st.session_state.upload_step = 1
    
    # Mostrar paso actual
    steps = {
        1: "📝 Metadata",
        2: "📄 Archivo",
        3: "⚙️ Procesamiento"
    }
    
    st.progress(st.session_state.upload_step / len(steps))
    st.subheader(f"Paso {st.session_state.upload_step}: {steps[st.session_state.upload_step]}")
    
    # Paso 1: Metadata
    if st.session_state.upload_step == 1:
        with st.form("metadata_form"):
            # Campos básicos
            title = st.text_input(
                "Título del Documento",
                help="Nombre descriptivo del documento"
            )
            
            # Categoría
            categories = doc_manager.categories["categories"]
            category = st.selectbox(
                "Categoría",
                options=list(categories.keys())
            )
            
            # Tipo y nivel
            col1, col2 = st.columns(2)
            with col1:
                doc_type = st.selectbox(
                    "Tipo de Documento",
                    options=doc_manager.get_document_types()
                )
            
            with col2:
                level = st.selectbox(
                    "Nivel",
                    options=doc_manager.get_difficulty_levels()
                )
            
            # Campos adicionales
            col3, col4 = st.columns(2)
            with col3:
                language = st.selectbox(
                    "Idioma",
                    options=["Español", "Inglés", "Francés", "Alemán"]
                )
                
                author = st.text_input(
                    "Autor",
                    help="Autor o creador del documento"
                )
            
            with col4:
                year = st.number_input(
                    "Año de Publicación",
                    min_value=1900,
                    max_value=2024,
                    value=2024
                )
                
                tags = st.text_input(
                    "Etiquetas",
                    help="Palabras clave separadas por comas"
                )
            
            description = st.text_area(
                "Descripción",
                help="Breve descripción del contenido del documento"
            )
            
            submitted = st.form_submit_button("Continuar")
        
        if submitted:
            if title:
                st.session_state.doc_metadata = {
                    "title": title,
                    "category": category,
                    "type": doc_type,
                    "level": level,
                    "language": language,
                    "author": author,
                    "year": year,
                    "tags": [tag.strip() for tag in tags.split(",") if tag.strip()],
                    "description": description
                }
                st.session_state.upload_step = 2
                st.rerun()
            else:
                st.error("❌ Por favor, completa al menos el título del documento.")
    
    # Paso 2: Subir archivo
    elif st.session_state.upload_step == 2:
        st.info("📝 Metadata configurada:")
        with st.expander("Ver detalles"):
            st.json(st.session_state.doc_metadata)
        
        uploaded_file = st.file_uploader(
            "Selecciona el documento",
            type=list(SUPPORTED_FORMATS.keys()),
            help="El documento será procesado, limpiado y convertido a vectorstore"
        )
        
        if uploaded_file:
            file_type = Path(uploaded_file.name).suffix.lower()[1:]
            st.info(f"""
            📄 Archivo seleccionado: {uploaded_file.name}
            - Tipo: {SUPPORTED_FORMATS[file_type][0]}
            - Tamaño: {uploaded_file.size / 1024:.1f} KB
            """)
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("← Volver", use_container_width=True):
                st.session_state.upload_step = 1
                st.rerun()
        
        with col2:
            if uploaded_file and st.button("Procesar →", use_container_width=True):
                st.session_state.uploaded_file = uploaded_file
                st.session_state.upload_step = 3
                st.rerun()
    
    # Paso 3: Procesamiento
    elif st.session_state.upload_step == 3:
        if hasattr(st.session_state, 'uploaded_file'):
            with st.spinner("⚙️ Procesando documento..."):
                with tempfile.TemporaryDirectory() as temp_dir:
                    result = process_document(
                        st.session_state.uploaded_file,
                        st.session_state.doc_metadata,
                        temp_dir
                    )
                    
                    if result["success"]:
                        try:
                            doc_hash = doc_manager.add_document(
                                st.session_state.doc_metadata,
                                result["vectorstore_path"],
                                result["original_path"]
                            )
                            
                            st.success(f"""
                            ✅ Documento procesado exitosamente:
                            - 📄 {result["num_pages"]} páginas procesadas
                            - 📚 {result["num_chunks"]} fragmentos generados
                            - 💾 {result["file_size"] / 1024:.1f} KB guardados
                            """)
                            
                            # Mostrar información y descargas
                            st.markdown("### 📑 Archivos Generados")
                            col1, col2 = st.columns(2)
                            
                            with col1:
                                st.markdown("**📁 Ubicación de archivos:**")
                                st.code(f"""
Documento original: {result['original_path']}
Vectorstore: {result['vectorstore_path']}
                                """)
                                
                                st.markdown("**💾 Descargas disponibles:**")
                                st.markdown(create_download_link(
                                    result['original_path'],
                                    "📥 Descargar documento original"
                                ), unsafe_allow_html=True)
                            with col2:
                                if result.get('preview_path'):
                                    st.image(
                                        result['preview_path'],
                                        caption="Vista previa del documento",
                                        use_column_width=True
                                    )
                                else:
                                    st.info("Vista previa no disponible para este formato")
                            
                            # Información del procesamiento
                            with st.expander("📊 Detalles del Procesamiento"):
                                st.markdown(f"""
                                **Información del documento:**
                                - Formato: {SUPPORTED_FORMATS[result['file_type']][0]}
                                - Páginas: {result['num_pages']}
                                - Fragmentos generados: {result['num_chunks']}
                                - Tamaño: {result['file_size'] / 1024:.1f} KB
                                
                                **Rutas del sistema:**
                                ```
                                {result['vectorstore_path']}
                                ```
                                
                                **Estado del procesamiento:**
                                - ✅ Documento original guardado
                                - ✅ Vectorstore generado
                                - {'✅' if result.get('preview_path') else '❌'} Vista previa generada
                                """)
                            
                            # Opciones post-procesamiento
                            st.markdown("### 🔄 Opciones")
                            col3, col4, col5 = st.columns(3)
                            
                            with col3:
                                if st.button("📤 Subir otro documento", use_container_width=True):
                                    # Limpiar session state
                                    for key in ['doc_metadata', 'uploaded_file']:
                                        if key in st.session_state:
                                            del st.session_state[key]
                                    st.session_state.upload_step = 1
                                    st.rerun()
                            
                            with col4:
                                if st.button("📚 Ir al Catálogo", use_container_width=True):
                                    st.switch_page("pages/1_📚_catalog.py")
                            
                            with col5:
                                if st.button("🤖 Crear Asistente", use_container_width=True):
                                    st.session_state.selected_docs = [doc_hash]
                                    st.switch_page("pages/2_🤖_agents.py")
                        
                        except Exception as e:
                            st.error(f"❌ Error al guardar el documento: {str(e)}")
                    else:
                        st.error(f"❌ Error al procesar el documento: {result['error']}")
                        if st.button("← Volver"):
                            st.session_state.upload_step = 2
                            st.rerun()

# Agregar estilos CSS personalizados
st.markdown("""
<style>
    .download-link {
        display: inline-block;
        padding: 0.5em 1em;
        background-color: #4CAF50;
        color: white;
        text-decoration: none;
        border-radius: 4px;
        transition: background-color 0.3s;
    }
    .download-link:hover {
        background-color: #45a049;
    }
    
    /* Estilo para los contenedores de documentos */
    .stExpander {
        border: 1px solid #ddd;
        border-radius: 8px;
        margin-bottom: 1rem;
    }
    
    /* Estilo para las pestañas de proceso */
    .stProgress > div > div > div {
        background-color: #4CAF50;
    }
    
    /* Estilo para los mensajes de éxito */
    .stSuccess {
        background-color: #e8f5e9;
        padding: 1rem;
        border-radius: 8px;
    }
</style>
""", unsafe_allow_html=True)

if __name__ == "__main__":
    main()
